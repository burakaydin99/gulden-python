!pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

slides_data = [
    ("Sustainable Cities and Communities", "Building the Future Together\nPrepared by: "),
    ("Why Sustainable Cities?", "• By 2050, 68% will live in cities\n• 75% of energy used by cities\n• 70% of emissions from cities\nSustainable cities reduce impact and improve life."),
    ("UN SDG 11", "Goal 11: Inclusive, safe, resilient, sustainable cities\n– Affordable housing and transport\n– Urban planning\n– Protect heritage\n– Expand green spaces"),
    ("Environmental Sustainability", "Renewable energy, green buildings, recycling, tree expansion"),
    ("Social Sustainability", "Equal access to services, inclusion, safe public spaces, culture preservation"),
    ("Economic Sustainability", "Green jobs, circular economy, innovation, sustainable transport"),
    ("Successful Cities Worldwide", "Copenhagen, Curitiba, Freiburg – sustainability works!"),
    ("Examples from Turkey", "Eskişehir – green parks\nIstanbul – smart city\nKonya – solar projects"),
    ("Challenges and Barriers", "Unplanned urbanization, traffic, pollution, inequality, climate risks"),
    ("What We Can Do", "Invest in green energy, promote sustainable transport, smart technologies, education"),
    ("Conclusion", "Sustainable cities are a necessity.\nTogether we can build cities respecting people and nature."),
    ("References", "UN SDG 11 • World Bank • Turkish Environment Agency • Greenpeace • ICLEI")
]

prs = Presentation()

for title, content in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(1))
    p = title_box.text_frame
    p.text = title
    p.paragraphs[0].font.bold = True
    p.paragraphs[0].font.size = Pt(28)
    p.paragraphs[0].font.color.rgb = RGBColor(76, 175, 80)
    
    body = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4))
    tf = body.text_frame
    tf.text = content
    for para in tf.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(51, 51, 51)

prs.save("Sustainable_Cities_and_Communities.pptx")
print("✅ PowerPoint created successfully!")
